from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
NAVY = RGBColor(0x0B, 0x1D, 0x51)
DARK_BLUE = RGBColor(0x1A, 0x3C, 0x7E)
MID_BLUE = RGBColor(0x2E, 0x5C, 0xAE)
LIGHT_BLUE = RGBColor(0x3A, 0x7B, 0xD5)
ACCENT_BLUE = RGBColor(0x5B, 0x9B, 0xD5)
SKY_BLUE = RGBColor(0xD6, 0xE8, 0xF7)
PALE_BLUE = RGBColor(0xEA, 0xF1, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
ORANGE = RGBColor(0xE8, 0x6C, 0x00)
TEAL = RGBColor(0x00, 0x89, 0x9B)
GREEN = RGBColor(0x2E, 0x7D, 0x32)


def add_bg_shape(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, font_name='맑은 고딕'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_text(text_frame, text, font_size=12, bold=False, color=DARK_GRAY, level=0, font_name='맑은 고딕', space_before=Pt(4), bullet_char='•'):
    p = text_frame.add_paragraph()
    if bullet_char:
        p.text = f"{bullet_char} {text}"
    else:
        p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.level = level
    p.space_before = space_before
    p.alignment = PP_ALIGN.LEFT
    return p


# ============================================================
# SLIDE 1: Title Slide
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# Background
add_bg_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), NAVY)

# Accent bar
add_bg_shape(slide, Inches(0), Inches(3.0), Inches(13.333), Inches(0.08), ACCENT_BLUE)

# Top label
add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(0.5),
             '한국기계연구원 (KIMM)', font_size=18, bold=False, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

# Main title
add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
             '3차년도 연구개발 계획', font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide, Inches(1.5), Inches(3.4), Inches(10), Inches(0.8),
             '무동력/비가압 담수화 시스템 성능 검증 및 담수화 설계 AI Agent 초기 모델 개발',
             font_size=20, bold=False, color=RGBColor(0xBB, 0xCC, 0xEE), alignment=PP_ALIGN.CENTER)

# Stage info
stage_box = add_text_box(slide, Inches(4.5), Inches(4.8), Inches(4.3), Inches(0.5),
                          '1단계  |  3차년도', font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Bottom line
add_bg_shape(slide, Inches(3), Inches(5.6), Inches(7.333), Inches(0.03), ACCENT_BLUE)

# Project name
add_text_box(slide, Inches(1.5), Inches(5.9), Inches(10), Inches(0.6),
             'AI 기반 무동력/비가압 태양에너지 활용 담수화 및 에너지 공급 시스템 설계 기술 개발',
             font_size=14, bold=False, color=RGBColor(0x88, 0x99, 0xBB), alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2: Overview / 연구개발 목표 및 연구내용 구성
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Header bar
add_bg_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), NAVY)
add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
             '3차년도 연구개발 목표 및 연구내용 구성', font_size=28, bold=True, color=WHITE)

# Goal box
goal_bg = add_bg_shape(slide, Inches(0.5), Inches(1.3), Inches(12.333), Inches(0.9), SKY_BLUE)
add_text_box(slide, Inches(0.7), Inches(1.35), Inches(1.5), Inches(0.8),
             '연구개발 목표', font_size=14, bold=True, color=DARK_BLUE)
add_text_box(slide, Inches(2.3), Inches(1.35), Inches(10), Inches(0.8),
             '무동력/비가압 담수화 시스템 성능 검증 및 담수화 설계 AI Agent 초기 모델 개발',
             font_size=16, bold=True, color=NAVY)

# 4 research content boxes
box_data = [
    ("연구내용 1", "연속형 담수화 시스템\nLab-scale 구축 및 성능 검증", MID_BLUE),
    ("연구내용 2", "담수화 시스템 운전 데이터\n수집 및 종합 성능 평가", LIGHT_BLUE),
    ("연구내용 3", "담수화 설계 AI Agent 1.0\n개발 및 워크플로우 구축", TEAL),
    ("연구내용 4", "(연구내용 3에 통합)", MED_GRAY),
]

x_start = 0.5
box_w = 2.85
gap = 0.2
y_top = 2.6

for i, (label, desc, color) in enumerate(box_data[:3]):
    x = x_start + i * (box_w + gap)
    # Number circle area
    add_bg_shape(slide, Inches(x), Inches(y_top), Inches(box_w), Inches(0.55), color)
    add_text_box(slide, Inches(x), Inches(y_top + 0.05), Inches(box_w), Inches(0.45),
                 label, font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    # Description box
    add_bg_shape(slide, Inches(x), Inches(y_top + 0.55), Inches(box_w), Inches(1.0), PALE_BLUE)
    add_text_box(slide, Inches(x + 0.15), Inches(y_top + 0.6), Inches(box_w - 0.3), Inches(0.9),
                 desc, font_size=13, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# Key items summary
y_summary = 4.5
add_bg_shape(slide, Inches(0.5), Inches(y_summary), Inches(12.333), Inches(0.04), ACCENT_BLUE)

items = [
    ("시스템 구축/검증", "히트펌프-담수화 모듈 Lab-scale 구축\n전기 인프라 유무·일사량 조건별 성능 실험"),
    ("데이터 수집/분석", "에너지자립형·연속형 성능 비교 분석\nAI Agent 학습용 데이터셋 구축"),
    ("AI Agent 개발", "담수화 설계 AI Agent 1.0 개발\n온톨로지/지식그래프, MCP 서버 구현\nAgent 간 워크플로우 엔진 구축"),
]

y_item = y_summary + 0.3
for i, (title, detail) in enumerate(items):
    x = 0.7 + i * 4.1
    add_text_box(slide, Inches(x), Inches(y_item), Inches(3.8), Inches(0.4),
                 title, font_size=14, bold=True, color=DARK_BLUE)
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y_item + 0.4), Inches(3.8), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for line in detail.split('\n'):
        p = tf.add_paragraph()
        p.text = f"• {line}"
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GRAY
        p.font.name = '맑은 고딕'
        p.space_before = Pt(2)


# ============================================================
# SLIDE 3: 연구내용 1 - 연속형 담수화 시스템 Lab-scale 구축 및 성능 검증
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_bg_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), NAVY)
add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
             '연구내용 1. 연속형 담수화 시스템 Lab-scale 구축 및 성능 검증', font_size=26, bold=True, color=WHITE)

# 1.1 Section
y = 1.3
add_bg_shape(slide, Inches(0.5), Inches(y), Inches(12.333), Inches(0.55), MID_BLUE)
add_text_box(slide, Inches(0.7), Inches(y + 0.05), Inches(11), Inches(0.45),
             '1.1 히트펌프-담수화 시스템 모듈 연계 Lab-scale 실험장치 구축', font_size=16, bold=True, color=WHITE)

y += 0.7
txBox = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(11.5), Inches(1.0))
tf = txBox.text_frame
tf.word_wrap = True
items_11 = [
    "히트펌프와 다중효용 담수화 시스템을 연계한 Lab-scale 실험장치 설계 및 제작",
    "증발기-응축기-히트펌프 간 열적 연계 최적화",
    "실시간 데이터 계측 및 수집 시스템 구축",
]
for item in items_11:
    p = tf.add_paragraph()
    p.text = f"• {item}"
    p.font.size = Pt(13)
    p.font.color.rgb = DARK_GRAY
    p.font.name = '맑은 고딕'
    p.space_before = Pt(4)

# 1.2 Section
y = 3.2
add_bg_shape(slide, Inches(0.5), Inches(y), Inches(12.333), Inches(0.55), MID_BLUE)
add_text_box(slide, Inches(0.7), Inches(y + 0.05), Inches(11), Inches(0.45),
             '1.2 연속형 담수화 시스템 성능 실험 및 평가', font_size=16, bold=True, color=WHITE)

y += 0.7
# Three columns for conditions
conditions = [
    ("전기 인프라 유무 조건별", "전기 인프라 유무 조건에 따른\n담수 생산 성능 비교 평가", "⚡"),
    ("일사량 조건별", "계절·시간대별 일사량 변화에\n따른 담수 생산 성능 평가", "☀"),
    ("히트펌프 운전 조건별", "히트펌프 운전 조건에 따른\n열전달 성능 분석", "🔄"),
]

for i, (title, desc, icon) in enumerate(conditions):
    x = 0.7 + i * 4.1
    add_bg_shape(slide, Inches(x), Inches(y), Inches(3.7), Inches(2.2), PALE_BLUE)
    # accent top bar
    add_bg_shape(slide, Inches(x), Inches(y), Inches(3.7), Inches(0.08), ACCENT_BLUE)
    add_text_box(slide, Inches(x + 0.2), Inches(y + 0.2), Inches(3.3), Inches(0.4),
                 title, font_size=14, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(x + 0.2), Inches(y + 0.7), Inches(3.3), Inches(1.3),
                 desc, font_size=12, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 4: 연구내용 2 - 담수화 시스템 운전 데이터 수집 및 종합 성능 평가
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_bg_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), NAVY)
add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
             '연구내용 2. 담수화 시스템 운전 데이터 수집 및 종합 성능 평가', font_size=26, bold=True, color=WHITE)

# 2.1
y = 1.3
add_bg_shape(slide, Inches(0.5), Inches(y), Inches(5.9), Inches(4.5), PALE_BLUE)
add_bg_shape(slide, Inches(0.5), Inches(y), Inches(5.9), Inches(0.08), MID_BLUE)
add_text_box(slide, Inches(0.7), Inches(y + 0.2), Inches(5.5), Inches(0.5),
             '2.1 에너지자립형·연속형 담수 생산 성능 비교·분석', font_size=16, bold=True, color=DARK_BLUE)

txBox = slide.shapes.add_textbox(Inches(0.9), Inches(y + 0.8), Inches(5.3), Inches(3.5))
tf = txBox.text_frame
tf.word_wrap = True
items_21 = [
    "에너지자립형 담수화 시스템 vs 연속형 담수화 시스템 성능 비교",
    "담수 생산량, 에너지 소비량, 수질 등 핵심 성능 지표 분석",
    "운전 조건별 최적 모드 전환 전략 도출",
    "시스템 효율 및 경제성 비교 분석",
]
for item in items_21:
    p = tf.add_paragraph()
    p.text = f"• {item}"
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GRAY
    p.font.name = '맑은 고딕'
    p.space_before = Pt(6)

# 2.2
add_bg_shape(slide, Inches(6.9), Inches(y), Inches(5.9), Inches(4.5), PALE_BLUE)
add_bg_shape(slide, Inches(6.9), Inches(y), Inches(5.9), Inches(0.08), TEAL)
add_text_box(slide, Inches(7.1), Inches(y + 0.2), Inches(5.5), Inches(0.5),
             '2.2 AI 설계 모델 개발을 위한 데이터 수집', font_size=16, bold=True, color=DARK_BLUE)

txBox = slide.shapes.add_textbox(Inches(7.3), Inches(y + 0.8), Inches(5.3), Inches(3.5))
tf = txBox.text_frame
tf.word_wrap = True
items_22 = [
    "담수화 시스템 운전 데이터 수집 체계 구축",
    "전기·열 수요 데이터 수집 및 분석",
    "수집 데이터의 AI Agent 학습용 데이터셋 변환",
    "데이터 품질 관리 및 정제 프로세스 확립",
    "학습 데이터 표준 포맷 정의 및 라벨링",
]
for item in items_22:
    p = tf.add_paragraph()
    p.text = f"• {item}"
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GRAY
    p.font.name = '맑은 고딕'
    p.space_before = Pt(6)

# Arrow / connection between the two
add_bg_shape(slide, Inches(5.5), Inches(3.2), Inches(2.2), Inches(0.5), WHITE)
add_text_box(slide, Inches(5.5), Inches(3.2), Inches(2.2), Inches(0.5),
             '성능 데이터 → AI 학습', font_size=11, bold=True, color=MID_BLUE, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 5: 연구내용 3 (Part 1) - 온톨로지/지식그래프 & MCP 서버
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_bg_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), NAVY)
add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
             '연구내용 3. 담수화 설계 AI Agent 1.0 개발 (1/2)', font_size=26, bold=True, color=WHITE)

# 3.1 Ontology
y = 1.3
add_bg_shape(slide, Inches(0.5), Inches(y), Inches(6.1), Inches(5.2), PALE_BLUE)
add_bg_shape(slide, Inches(0.5), Inches(y), Inches(6.1), Inches(0.55), TEAL)
add_text_box(slide, Inches(0.7), Inches(y + 0.05), Inches(5.7), Inches(0.45),
             '3.1 담수화 도메인 온톨로지/지식그래프 1.0 구축', font_size=15, bold=True, color=WHITE)

txBox = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.7), Inches(5.7), Inches(4.3))
tf = txBox.text_frame
tf.word_wrap = True
items_31 = [
    ("시스템 레벨 온톨로지 확장", "담수화 모듈-에너지 공급-전처리 간\n인터페이스 관계 모델링"),
    ("온톨로지 기반 설계 추론 엔진 개발", "제약조건 전파 및 설계 일관성 검증"),
    ("지식 자동 업데이트 체계", "신규 문헌/실험 데이터 자동 추출 및\n지식 업데이트 체계 구축"),
]
for title, desc in items_31:
    p = tf.add_paragraph()
    p.text = f"▶ {title}"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = '맑은 고딕'
    p.space_before = Pt(10)

    for line in desc.split('\n'):
        p2 = tf.add_paragraph()
        p2.text = f"   - {line}"
        p2.font.size = Pt(11)
        p2.font.color.rgb = DARK_GRAY
        p2.font.name = '맑은 고딕'
        p2.space_before = Pt(2)

# 3.2 MCP Server
add_bg_shape(slide, Inches(6.9), Inches(y), Inches(5.9), Inches(5.2), PALE_BLUE)
add_bg_shape(slide, Inches(6.9), Inches(y), Inches(5.9), Inches(0.55), MID_BLUE)
add_text_box(slide, Inches(7.1), Inches(y + 0.05), Inches(5.5), Inches(0.45),
             '3.2 담수화 설계 S/W API 및 MCP 서버 구현', font_size=15, bold=True, color=WHITE)

txBox = slide.shapes.add_textbox(Inches(7.1), Inches(y + 0.7), Inches(5.5), Inches(4.3))
tf = txBox.text_frame
tf.word_wrap = True
items_32 = [
    ("설계 모듈별 RESTful API 개발", "MED 설계, 히트펌프 사이클,\n열교환기 최적화"),
    ("AI Agent 전용 MCP 서버 구현", "Tool 스키마 표준화,\n시뮬레이션 SW 래핑"),
    ("외부 AI 모델 연동", "서울시립대 담수량 예측 모델 등\n연동 인터페이스 설계"),
]
for title, desc in items_32:
    p = tf.add_paragraph()
    p.text = f"▶ {title}"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = '맑은 고딕'
    p.space_before = Pt(10)

    for line in desc.split('\n'):
        p2 = tf.add_paragraph()
        p2.text = f"   - {line}"
        p2.font.size = Pt(11)
        p2.font.color.rgb = DARK_GRAY
        p2.font.name = '맑은 고딕'
        p2.space_before = Pt(2)


# ============================================================
# SLIDE 6: 연구내용 3 (Part 2) - AI Agent 1.0 & 워크플로우
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_bg_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), NAVY)
add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
             '연구내용 3. 담수화 설계 AI Agent 1.0 개발 (2/2)', font_size=26, bold=True, color=WHITE)

# 3.3 AI Agent 1.0
y = 1.3
add_bg_shape(slide, Inches(0.5), Inches(y), Inches(12.333), Inches(0.55), DARK_BLUE)
add_text_box(slide, Inches(0.7), Inches(y + 0.05), Inches(11), Inches(0.45),
             '3.3 담수화 설계 AI Agent 1.0 개발', font_size=16, bold=True, color=WHITE)

# Agent cards
agents = [
    ("증류기 구조/형상\n설계 Agent",
     ["자연어 입력 해석", "설계 파라미터 자동 결정", "검증/피드백 루프"],
     MID_BLUE),
    ("히트펌프-열교환기\n설계 Agent",
     ["사이클 구성", "열전달 유로 형상 탐색", "CFD 연동 검증"],
     LIGHT_BLUE),
    ("에너지 공급 시스템\n설계 Agent",
     ["전기·열 부하 분석", "재생에너지원 조합 최적화"],
     TEAL),
    ("HITL 설계 검토\n프로세스",
     ["Human-In-The-Loop", "Agent 핵심 기능 검증", "설계 검토 확립"],
     GREEN),
]

y_agent = y + 0.7
for i, (title, items, color) in enumerate(agents):
    x = 0.5 + i * 3.15
    # Card
    add_bg_shape(slide, Inches(x), Inches(y_agent), Inches(2.95), Inches(2.8), LIGHT_GRAY)
    # Color top bar
    add_bg_shape(slide, Inches(x), Inches(y_agent), Inches(2.95), Inches(0.08), color)
    # Title
    add_text_box(slide, Inches(x + 0.15), Inches(y_agent + 0.2), Inches(2.65), Inches(0.7),
                 title, font_size=13, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

    txBox = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y_agent + 0.95), Inches(2.55), Inches(1.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    for item in items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GRAY
        p.font.name = '맑은 고딕'
        p.space_before = Pt(3)

# 3.4 Workflow Engine
y_wf = y_agent + 3.1
add_bg_shape(slide, Inches(0.5), Inches(y_wf), Inches(12.333), Inches(0.55), DARK_BLUE)
add_text_box(slide, Inches(0.7), Inches(y_wf + 0.05), Inches(11), Inches(0.45),
             '3.4 AI Agent 간 워크플로우 엔진 구축', font_size=16, bold=True, color=WHITE)

y_wf_items = y_wf + 0.7
wf_items = [
    ("Agent 간 통신 프로토콜", "Agent 간 통신 프로토콜 및\n작업 위임 체계 구현"),
    ("설계 워크플로우 시각화", "설계 워크플로우\n시각화 도구 개발"),
    ("빌더-Agent 통합", "노코드 빌더에서\n설계 Agent 실시간 호출"),
]

for i, (title, desc) in enumerate(wf_items):
    x = 0.7 + i * 4.1
    add_bg_shape(slide, Inches(x), Inches(y_wf_items), Inches(3.7), Inches(1.2), PALE_BLUE)
    add_bg_shape(slide, Inches(x), Inches(y_wf_items), Inches(3.7), Inches(0.06), ACCENT_BLUE)
    add_text_box(slide, Inches(x + 0.15), Inches(y_wf_items + 0.1), Inches(3.4), Inches(0.35),
                 title, font_size=13, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(x + 0.15), Inches(y_wf_items + 0.45), Inches(3.4), Inches(0.7),
                 desc, font_size=11, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 7: Summary / 3차년도 핵심 성과 목표
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_bg_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), NAVY)
add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
             '3차년도 핵심 성과 목표 요약', font_size=28, bold=True, color=WHITE)

# Three major goals
goals = [
    {
        "title": "시스템 구축 및 검증",
        "color": MID_BLUE,
        "items": [
            "히트펌프-담수화 모듈 연계 Lab-scale 실험장치 구축",
            "전기 인프라 유무 조건별 성능 비교 평가",
            "일사량 조건별 담수 생산 성능 평가",
            "에너지자립형·연속형 성능 비교 분석",
            "AI Agent 학습용 운전 데이터 수집",
        ]
    },
    {
        "title": "AI 인프라 구축",
        "color": TEAL,
        "items": [
            "담수화 도메인 온톨로지/지식그래프 1.0",
            "설계 S/W RESTful API 개발",
            "AI Agent 전용 MCP 서버 구현",
            "외부 AI 모델 연동 인터페이스",
        ]
    },
    {
        "title": "AI Agent 개발",
        "color": ORANGE,
        "items": [
            "증류기 설계 / 히트펌프 설계 / 에너지 설계 Agent",
            "HITL 설계 검토 프로세스 확립",
            "Agent 간 워크플로우 엔진 구축",
            "노코드 빌더-Agent 통합",
        ]
    },
]

for i, goal in enumerate(goals):
    x = 0.5 + i * 4.2
    y = 1.3

    # Header
    add_bg_shape(slide, Inches(x), Inches(y), Inches(3.9), Inches(0.6), goal["color"])
    add_text_box(slide, Inches(x), Inches(y + 0.07), Inches(3.9), Inches(0.45),
                 goal["title"], font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Items
    add_bg_shape(slide, Inches(x), Inches(y + 0.6), Inches(3.9), Inches(4.5), LIGHT_GRAY)
    txBox = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.8), Inches(3.5), Inches(4.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    for item in goal["items"]:
        p = tf.add_paragraph()
        p.text = f"✓ {item}"
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
        p.font.name = '맑은 고딕'
        p.space_before = Pt(8)


# Save
output_path = r"C:\Users\USER\Documents\퇴직준비자_임병주\AI담수\KIMM_3차년도_연구계획.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
